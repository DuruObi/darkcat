from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

def add_slide(title, content, title_color=(0,0,0), content_color=(0,0,0)):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color)
    tf = slide.placeholders[1].text_frame
    tf.text = content
    tf.paragraphs[0].font.color.rgb = RGBColor(*content_color)

# Slide 1 — Cover Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "DarkCat 🐈‍⬛"
slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(60)
slide.placeholders[1].text = "Developer Automation CLI\nOEDX DIGITAL"

# Add badges to cover
left = Inches(1)
top = Inches(3.5)
slide.shapes.add_picture("pypi-badge.png", left, top, width=Inches(2.5))
slide.shapes.add_picture("ci-badge.png", left + Inches(3), top, width=Inches(2.5))

# Styled slides
text_color = (30, 30, 30)
title_color = (0, 80, 160)

add_slide(
    "What is DarkCat?",
    "• Terminal-based developer automation tool\n"
    "• Scaffolds projects, adds automation, and runs system checks\n"
    "• Designed for speed, modularity, and scalability",
    title_color, text_color
)

add_slide(
    "Key Features",
    "• Interactive wizard mode\n"
    "• Dynamic templates (web & API)\n"
    "• Docker & .env automation\n"
    "• System diagnostics (doctor)\n"
    "• Config file support (~/.darkcat.yaml)\n"
    "• Plugin system for extensions\n"
    "• Versioned & PyPI installable",
    title_color, text_color
)

add_slide(
    "How DarkCat Works",
    "User runs CLI → Wizard / Init / Add / Doctor → Project scaffolded → Optional plugins loaded",
    title_color, text_color
)

add_slide(
    "Commands Overview",
    "init       → Scaffold new projects\n"
    "wizard     → Guided interactive setup\n"
    "add        → Add Docker/env support\n"
    "doctor     → System diagnostics\n"
    "--version  → Show CLI version\n"
    "Plugins    → Extend CLI with external commands",
    title_color, text_color
)

add_slide(
    "Installation",
    "pip install darkcat\nTest:\ndarkcat --version\ndarkcat wizard",
    title_color, text_color
)

add_slide(
    "Config File Support",
    "File: ~/.darkcat.yaml\nExample:\n"
    "default_template: default\n"
    "last_project_name: MyWebApp\n"
    "default_folder: ~/Projects",
    title_color, text_color
)

add_slide(
    "Plugin System",
    "• Drop Python files into darkcat/plugins/\n"
    "• Each plugin defines a cli_command\n"
    "• Automatically registered with CLI\n"
    "Example:\n@click.command()\ndef cli_command():\n    click.echo('🐈‍⬛ Plugin works!')",
    title_color, text_color
)

add_slide(
    "Key Benefits",
    "• Fast CLI workflow\n• Modular & extendable\n• CI/CD ready & PyPI installable\n"
    "• Configurable via YAML\n• Portfolio-friendly & repeatable",
    title_color, text_color
)

add_slide(
    "Step-by-Step Workflow",
    "1. darkcat --version\n2. darkcat wizard\n3. darkcat add docker --folder MyWebApp\n"
    "4. darkcat add env --folder MyWebApp\n5. darkcat doctor\n6. darkcat sample-plugin",
    title_color, text_color
)

add_slide(
    "Roadmap / Next Steps",
    "• DarkCat Pro features\n• More templates & themes\n• Auto-update checker\n"
    "• Enhanced plugin ecosystem\n• Web dashboard integration",
    title_color, text_color
)

add_slide(
    "Contact / GitHub",
    "GitHub: github.com/OEDX-DIGITAL/darkcat\nPyPI: pypi.org/project/darkcat\nQuestions? Feedback welcome!",
    title_color, text_color
)

# Save presentation
prs.save("darkcat_presentation_styled.pptx")
print("✅ Styled DarkCat presentation created: darkcat_presentation_styled.pptx")
